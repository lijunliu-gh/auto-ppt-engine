"""Tests for python_backend.visual_qa module.

Covers Issue, helpers, analyze_visual_quality, _export_images, and run_visual_qa.
"""

from __future__ import annotations

import json
import subprocess
from pathlib import Path
from types import SimpleNamespace
from typing import Any
from unittest.mock import MagicMock, patch

import pytest

from python_backend.visual_qa import (
    EMU_PER_INCH,
    Issue,
    _boxes_overlap,
    _emu_to_inches,
    _export_images,
    _run_command,
    _shape_bbox_in_inches,
    analyze_visual_quality,
    run_visual_qa,
)


# ── Issue dataclass ──────────────────────────────────────────────────────────


class TestIssue:
    def test_default_severity(self):
        issue = Issue(code="test", message="msg")
        assert issue.severity == "warning"

    def test_custom_severity(self):
        issue = Issue(code="x", message="y", severity="high")
        assert issue.severity == "high"

    def test_to_dict(self):
        issue = Issue(code="edge_margin", message="Too close", severity="info")
        d = issue.to_dict()
        assert d == {"code": "edge_margin", "message": "Too close", "severity": "info"}

    def test_to_dict_keys(self):
        d = Issue(code="a", message="b").to_dict()
        assert set(d.keys()) == {"code", "message", "severity"}


# ── Helper functions ─────────────────────────────────────────────────────────


class TestEmuToInches:
    def test_zero(self):
        assert _emu_to_inches(0) == 0.0

    def test_one_inch(self):
        assert _emu_to_inches(EMU_PER_INCH) == pytest.approx(1.0)

    def test_half_inch(self):
        assert _emu_to_inches(EMU_PER_INCH // 2) == pytest.approx(0.5, abs=0.01)

    def test_negative(self):
        assert _emu_to_inches(-EMU_PER_INCH) == pytest.approx(-1.0)


class TestShapeBboxInInches:
    @staticmethod
    def _shape(left: int, top: int, width: int, height: int) -> SimpleNamespace:
        return SimpleNamespace(left=left, top=top, width=width, height=height)

    def test_normal_shape(self):
        s = self._shape(EMU_PER_INCH, EMU_PER_INCH, 2 * EMU_PER_INCH, 3 * EMU_PER_INCH)
        bbox = _shape_bbox_in_inches(s)
        assert bbox is not None
        x1, y1, x2, y2 = bbox
        assert x1 == pytest.approx(1.0)
        assert y1 == pytest.approx(1.0)
        assert x2 == pytest.approx(3.0)
        assert y2 == pytest.approx(4.0)

    def test_zero_width_returns_none(self):
        s = self._shape(0, 0, 0, EMU_PER_INCH)
        assert _shape_bbox_in_inches(s) is None

    def test_zero_height_returns_none(self):
        s = self._shape(0, 0, EMU_PER_INCH, 0)
        assert _shape_bbox_in_inches(s) is None

    def test_negative_width_returns_none(self):
        s = self._shape(0, 0, -1, EMU_PER_INCH)
        assert _shape_bbox_in_inches(s) is None

    def test_exception_returns_none(self):
        s = SimpleNamespace(left="bad", top=0, width=100, height=100)
        # int("bad") will raise — function should return None
        assert _shape_bbox_in_inches(s) is None

    def test_missing_attr_returns_none(self):
        assert _shape_bbox_in_inches(object()) is None


class TestBoxesOverlap:
    def test_overlapping(self):
        a = (0.0, 0.0, 2.0, 2.0)
        b = (1.0, 1.0, 3.0, 3.0)
        assert _boxes_overlap(a, b) is True

    def test_no_overlap_right(self):
        a = (0.0, 0.0, 1.0, 1.0)
        b = (2.0, 0.0, 3.0, 1.0)
        assert _boxes_overlap(a, b) is False

    def test_no_overlap_below(self):
        a = (0.0, 0.0, 1.0, 1.0)
        b = (0.0, 2.0, 1.0, 3.0)
        assert _boxes_overlap(a, b) is False

    def test_touching_edges_not_overlapping(self):
        a = (0.0, 0.0, 1.0, 1.0)
        b = (1.0, 0.0, 2.0, 1.0)
        assert _boxes_overlap(a, b) is False

    def test_contained(self):
        outer = (0.0, 0.0, 10.0, 10.0)
        inner = (2.0, 2.0, 4.0, 4.0)
        assert _boxes_overlap(outer, inner) is True

    def test_identical(self):
        a = (1.0, 1.0, 3.0, 3.0)
        assert _boxes_overlap(a, a) is True


# ── _run_command ─────────────────────────────────────────────────────────────


class TestRunCommand:
    def test_success(self):
        ok, err = _run_command(["echo", "hello"])
        assert ok is True
        assert err == ""

    def test_command_not_found(self):
        ok, err = _run_command(["__nonexistent_binary_xyz__"])
        assert ok is False
        assert "not found" in err.lower() or "Command not found" in err

    def test_command_failure(self):
        ok, err = _run_command(["false"])
        assert ok is False


# ── _export_images ───────────────────────────────────────────────────────────


class TestExportImages:
    def test_no_soffice(self, tmp_path: Path):
        with patch("python_backend.visual_qa.shutil.which", return_value=None):
            images, notes = _export_images(tmp_path / "deck.pptx", tmp_path / "imgs", 150)
        assert images == []
        assert any("soffice" in n for n in notes)

    def test_no_pdftoppm(self, tmp_path: Path):
        def which_side(cmd: str) -> str | None:
            return "/usr/bin/soffice" if cmd == "soffice" else None

        with patch("python_backend.visual_qa.shutil.which", side_effect=which_side):
            images, notes = _export_images(tmp_path / "deck.pptx", tmp_path / "imgs", 150)
        assert images == []
        assert any("pdftoppm" in n for n in notes)

    def test_soffice_conversion_fails(self, tmp_path: Path):
        with (
            patch("python_backend.visual_qa.shutil.which", return_value="/usr/bin/mock"),
            patch("python_backend.visual_qa._run_command", return_value=(False, "conversion error")),
        ):
            images, notes = _export_images(tmp_path / "deck.pptx", tmp_path / "imgs", 150)
        assert images == []
        assert any("conversion" in n.lower() for n in notes)

    def test_pdf_not_produced(self, tmp_path: Path):
        # soffice succeeds but no PDF appears
        with (
            patch("python_backend.visual_qa.shutil.which", return_value="/usr/bin/mock"),
            patch("python_backend.visual_qa._run_command", return_value=(True, "")),
        ):
            images, notes = _export_images(tmp_path / "deck.pptx", tmp_path / "imgs", 150)
        assert images == []
        assert any("no PDF" in n or "no pdf" in n.lower() for n in notes)

    def test_pdftoppm_fails(self, tmp_path: Path):
        imgs_dir = tmp_path / "imgs"
        imgs_dir.mkdir(parents=True, exist_ok=True)
        # Create fake PDF so the file-exists check passes
        (imgs_dir / "deck.pdf").write_text("fake")

        call_count = 0

        def fake_run_command(cmd: list[str]) -> tuple[bool, str]:
            nonlocal call_count
            call_count += 1
            if call_count == 1:
                return True, ""  # soffice succeeds
            return False, "pdftoppm failed"  # pdftoppm fails

        with (
            patch("python_backend.visual_qa.shutil.which", return_value="/usr/bin/mock"),
            patch("python_backend.visual_qa._run_command", side_effect=fake_run_command),
        ):
            images, notes = _export_images(tmp_path / "deck.pptx", imgs_dir, 150)
        assert images == []
        assert any("pdftoppm" in n.lower() or "image conversion" in n.lower() for n in notes)

    def test_no_images_generated(self, tmp_path: Path):
        imgs_dir = tmp_path / "imgs"
        imgs_dir.mkdir(parents=True, exist_ok=True)
        (imgs_dir / "deck.pdf").write_text("fake")

        with (
            patch("python_backend.visual_qa.shutil.which", return_value="/usr/bin/mock"),
            patch("python_backend.visual_qa._run_command", return_value=(True, "")),
        ):
            images, notes = _export_images(tmp_path / "deck.pptx", imgs_dir, 150)
        assert images == []
        assert any("no slide images" in n.lower() for n in notes)


# ── analyze_visual_quality ───────────────────────────────────────────────────


def _make_pptx(tmp_path: Path, shapes_per_slide: list[list[dict[str, Any]]] | None = None) -> Path:
    """Create a minimal PPTX for testing with configurable shapes."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    if shapes_per_slide is None:
        shapes_per_slide = [[]]

    for slide_shapes in shapes_per_slide:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        for shape_spec in slide_shapes:
            left = Inches(shape_spec.get("left", 1))
            top = Inches(shape_spec.get("top", 1))
            width = Inches(shape_spec.get("width", 2))
            height = Inches(shape_spec.get("height", 1))
            text = shape_spec.get("text", "")
            tf = slide.shapes.add_textbox(left, top, width, height)
            tf.text_frame.text = text

    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


class TestAnalyzeVisualQuality:
    def test_empty_slide_detected(self, tmp_path: Path):
        pptx = _make_pptx(tmp_path, [[]])
        result = analyze_visual_quality(pptx)
        assert result["slideCount"] == 1
        issues = result["slides"][0]["issues"]
        codes = [i["code"] for i in issues]
        assert "empty_slide" in codes

    def test_single_text_shape_no_empty(self, tmp_path: Path):
        pptx = _make_pptx(tmp_path, [[{"text": "Hello", "left": 1, "top": 1}]])
        result = analyze_visual_quality(pptx)
        issues = result["slides"][0]["issues"]
        codes = [i["code"] for i in issues]
        assert "empty_slide" not in codes

    def test_edge_margin_detected(self, tmp_path: Path):
        # Place shape at extreme left edge
        pptx = _make_pptx(tmp_path, [[{"text": "Edge", "left": 0.05, "top": 1, "width": 1, "height": 0.5}]])
        result = analyze_visual_quality(pptx, margin_in=0.3)
        issues = result["slides"][0]["issues"]
        codes = [i["code"] for i in issues]
        assert "edge_margin" in codes

    def test_no_edge_margin_when_centered(self, tmp_path: Path):
        pptx = _make_pptx(tmp_path, [[{"text": "Center", "left": 2, "top": 2, "width": 2, "height": 1}]])
        result = analyze_visual_quality(pptx, margin_in=0.3)
        issues = result["slides"][0]["issues"]
        edge_issues = [i for i in issues if i["code"] == "edge_margin"]
        assert len(edge_issues) == 0

    def test_overlap_detected(self, tmp_path: Path):
        shapes = [
            {"text": "A", "left": 1, "top": 1, "width": 3, "height": 2},
            {"text": "B", "left": 2, "top": 1.5, "width": 3, "height": 2},
        ]
        pptx = _make_pptx(tmp_path, [shapes])
        result = analyze_visual_quality(pptx)
        issues = result["slides"][0]["issues"]
        codes = [i["code"] for i in issues]
        assert "overlap_candidate" in codes

    def test_no_overlap_when_separate(self, tmp_path: Path):
        shapes = [
            {"text": "A", "left": 0.5, "top": 1, "width": 1, "height": 1},
            {"text": "B", "left": 5, "top": 1, "width": 1, "height": 1},
        ]
        pptx = _make_pptx(tmp_path, [shapes])
        result = analyze_visual_quality(pptx)
        issues = result["slides"][0]["issues"]
        overlap_issues = [i for i in issues if i["code"] == "overlap_candidate"]
        assert len(overlap_issues) == 0

    def test_text_only_for_middle_slide(self, tmp_path: Path):
        # 3 slides: first, middle (text-only), last
        shapes_per_slide = [
            [{"text": "Title", "left": 1, "top": 1}],
            [{"text": "Only text", "left": 1, "top": 1}],
            [{"text": "End", "left": 1, "top": 1}],
        ]
        pptx = _make_pptx(tmp_path, shapes_per_slide)
        result = analyze_visual_quality(pptx)
        assert result["slideCount"] == 3
        # Middle slide (index 1) should have text_only
        middle_issues = result["slides"][1]["issues"]
        codes = [i["code"] for i in middle_issues]
        assert "text_only" in codes

    def test_text_only_not_for_first_slide(self, tmp_path: Path):
        shapes_per_slide = [
            [{"text": "First", "left": 1, "top": 1}],
            [{"text": "Middle", "left": 1, "top": 1}],
            [{"text": "Last", "left": 1, "top": 1}],
        ]
        pptx = _make_pptx(tmp_path, shapes_per_slide)
        result = analyze_visual_quality(pptx)
        first_issues = result["slides"][0]["issues"]
        codes = [i["code"] for i in first_issues]
        assert "text_only" not in codes

    def test_text_only_not_for_last_slide(self, tmp_path: Path):
        shapes_per_slide = [
            [{"text": "First", "left": 1, "top": 1}],
            [{"text": "Middle", "left": 1, "top": 1}],
            [{"text": "Last", "left": 1, "top": 1}],
        ]
        pptx = _make_pptx(tmp_path, shapes_per_slide)
        result = analyze_visual_quality(pptx)
        last_issues = result["slides"][2]["issues"]
        codes = [i["code"] for i in last_issues]
        assert "text_only" not in codes

    def test_multiple_slides_reported(self, tmp_path: Path):
        shapes = [[{"text": f"S{i}", "left": 1, "top": 1}] for i in range(5)]
        pptx = _make_pptx(tmp_path, shapes)
        result = analyze_visual_quality(pptx)
        assert result["slideCount"] == 5
        assert len(result["slides"]) == 5
        for i, slide_report in enumerate(result["slides"]):
            assert slide_report["slide"] == i + 1


# ── run_visual_qa ────────────────────────────────────────────────────────────


class TestRunVisualQa:
    def test_file_not_found_raises(self, tmp_path: Path):
        with pytest.raises(RuntimeError, match="PPTX file not found"):
            run_visual_qa(tmp_path / "nonexistent.pptx")

    def test_full_report_structure(self, tmp_path: Path):
        pptx = _make_pptx(tmp_path, [[{"text": "Hello", "left": 1, "top": 1}]])
        output = tmp_path / "qa-out"
        report = run_visual_qa(pptx, output_dir=output)

        assert report["pptxPath"] == str(pptx.resolve())
        assert "generatedAt" in report
        assert "analysis" in report
        assert "summary" in report
        assert report["summary"]["slideCount"] == 1
        assert isinstance(report["summary"]["totalIssues"], int)
        assert isinstance(report["summary"]["highRiskSlides"], list)
        assert "reportPath" in report

        # Check that report JSON was written
        report_file = Path(report["reportPath"])
        assert report_file.exists()
        saved = json.loads(report_file.read_text())
        assert saved["summary"]["slideCount"] == 1

    def test_default_output_dir(self, tmp_path: Path):
        pptx = _make_pptx(tmp_path, [[{"text": "Hi", "left": 1, "top": 1}]])
        report = run_visual_qa(pptx)
        # Default dir should be <stem>-qa next to the pptx
        expected_dir = pptx.resolve().parent / f"{pptx.stem}-qa"
        assert Path(report["reportPath"]).parent == expected_dir
        # Clean up
        import shutil
        shutil.rmtree(expected_dir, ignore_errors=True)

    def test_high_risk_slides_detected(self, tmp_path: Path):
        # Empty slide => severity=high
        pptx = _make_pptx(tmp_path, [[]])
        report = run_visual_qa(pptx, output_dir=tmp_path / "qa")
        assert 1 in report["summary"]["highRiskSlides"]

    def test_notes_from_image_export(self, tmp_path: Path):
        pptx = _make_pptx(tmp_path, [[{"text": "X", "left": 1, "top": 1}]])
        # soffice not available → should gracefully add a note
        with patch("python_backend.visual_qa.shutil.which", return_value=None):
            report = run_visual_qa(pptx, output_dir=tmp_path / "qa")
        assert any("soffice" in n for n in report["notes"])
        assert report["images"] == []
