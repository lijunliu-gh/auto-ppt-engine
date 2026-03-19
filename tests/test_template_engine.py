"""Tests for template_engine and pptx_renderer modules."""

from __future__ import annotations

import json
import os
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation
from pptx.util import Inches

ROOT_DIR = Path(__file__).resolve().parent.parent

# ── Helpers ──────────────────────────────────────────────────────────


def _create_test_template(tmp_path: Path, num_layouts: int = 2) -> Path:
    """Create a minimal .pptx template file for testing."""
    prs = Presentation()
    # The default presentation already has slide layouts from the default master
    path = tmp_path / "test-template.pptx"
    prs.save(str(path))
    return path


def _build_mock_deck() -> dict:
    """Build a minimal valid deck JSON for testing."""
    return {
        "schemaVersion": "0.3.0",
        "deckTitle": "Test Deck",
        "language": "en",
        "audience": "Testers",
        "scenario": "Unit testing",
        "tone": "technical",
        "theme": "default",
        "sourceDisplayMode": "notes",
        "slideCount": 5,
        "needsSpeakerNotes": True,
        "assumptions": ["This is a test"],
        "slides": [
            {
                "page": 1, "layout": "title", "title": "Test Title",
                "subtitle": "Test Subtitle", "objective": "", "bullets": [],
                "left": [], "right": [], "table": {"columns": [], "rows": []},
                "chart": {"type": "bar", "title": "", "categories": [], "series": []},
                "visuals": [], "sources": [], "speakerNotes": ["Welcome"]
            },
            {
                "page": 2, "layout": "bullet", "title": "Key Points",
                "subtitle": "", "objective": "Understand key points",
                "bullets": ["Point A", "Point B", "Point C"],
                "left": [], "right": [], "table": {"columns": [], "rows": []},
                "chart": {"type": "bar", "title": "", "categories": [], "series": []},
                "visuals": ["A diagram showing flow"], "sources": [
                    {"id": "s1", "label": "Source 1", "type": "text", "location": "file.txt",
                     "trustLevel": "user-provided", "priority": "high", "usedFor": ["context"]}
                ],
                "speakerNotes": ["Discuss points"]
            },
            {
                "page": 3, "layout": "two-column", "title": "Comparison",
                "subtitle": "", "objective": "",
                "bullets": [],
                "left": ["Left A", "Left B"], "right": ["Right A", "Right B"],
                "table": {"columns": [], "rows": []},
                "chart": {"type": "bar", "title": "", "categories": [], "series": []},
                "visuals": [], "sources": [], "speakerNotes": []
            },
            {
                "page": 4, "layout": "chart", "title": "Revenue Chart",
                "subtitle": "", "objective": "",
                "bullets": ["Revenue up 20%"],
                "left": [], "right": [],
                "table": {"columns": [], "rows": []},
                "chart": {
                    "type": "bar", "title": "Q Revenue",
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [{"name": "2024", "data": [100, 150, 200, 250]}]
                },
                "visuals": [], "sources": [], "speakerNotes": []
            },
            {
                "page": 5, "layout": "closing", "title": "Thank You",
                "subtitle": "Questions?", "objective": "",
                "bullets": ["Contact: test@example.com"],
                "left": [], "right": [],
                "table": {"columns": [], "rows": []},
                "chart": {"type": "bar", "title": "", "categories": [], "series": []},
                "visuals": [], "sources": [], "speakerNotes": []
            },
        ]
    }


# ── Template Engine Tests ────────────────────────────────────────────


class TestTemplateEngine:
    def test_parse_template_returns_config(self, tmp_path):
        from python_backend.template_engine import parse_template
        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        assert config.source_path == tpl.resolve()
        assert len(config.layouts) > 0
        assert isinstance(config.layout_mapping, dict)

    def test_parse_template_extracts_layouts(self, tmp_path):
        from python_backend.template_engine import parse_template
        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        # Default pptx has multiple slide layouts
        assert len(config.layouts) >= 1
        for layout in config.layouts:
            assert layout.name  # Each layout has a name

    def test_parse_template_file_not_found(self, tmp_path):
        from python_backend.template_engine import parse_template
        with pytest.raises(FileNotFoundError):
            parse_template(tmp_path / "nonexistent.pptx")

    def test_parse_template_wrong_extension(self, tmp_path):
        from python_backend.template_engine import parse_template
        bad_file = tmp_path / "template.docx"
        bad_file.write_bytes(b"fake")
        with pytest.raises(ValueError, match=".pptx"):
            parse_template(bad_file)

    def test_parse_template_too_large(self, tmp_path):
        from python_backend.template_engine import parse_template, MAX_TEMPLATE_SIZE
        huge_file = tmp_path / "huge.pptx"
        huge_file.write_bytes(b"\0" * (MAX_TEMPLATE_SIZE + 1))
        with pytest.raises(ValueError, match="limit"):
            parse_template(huge_file)

    def test_layout_mapping_has_common_layouts(self, tmp_path):
        from python_backend.template_engine import parse_template
        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        # At minimum, the auto-mapper should map something
        assert len(config.layout_mapping) > 0

    def test_describe_template(self, tmp_path):
        from python_backend.template_engine import parse_template, describe_template
        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        desc = describe_template(config)
        assert "layoutCount" in desc
        assert "layouts" in desc
        assert "layoutMapping" in desc
        assert "fonts" in desc
        assert "themeColors" in desc

    def test_layout_info_properties(self, tmp_path):
        from python_backend.template_engine import parse_template
        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        # Check LayoutInfo property accessors
        for layout in config.layouts:
            assert isinstance(layout.has_title, bool)
            assert isinstance(layout.has_body, bool)
            assert isinstance(layout.has_picture, bool)
            assert isinstance(layout.body_count, int)

    def test_theme_colors_defaults(self):
        from python_backend.template_engine import ThemeColors
        colors = ThemeColors()
        assert colors.primary == "0F766E"
        assert colors.secondary == "2563EB"
        assert colors.background == "FFFFFF"

    def test_placeholder_info_fields(self):
        from python_backend.template_engine import PlaceholderInfo
        ph = PlaceholderInfo(idx=0, name="Title 1", placeholder_type="TITLE",
                            left=0, top=0, width=1000, height=500)
        assert ph.idx == 0
        assert ph.placeholder_type == "TITLE"


# ── PPTX Renderer Tests ─────────────────────────────────────────────


class TestPptxRenderer:
    def test_render_creates_pptx_file(self, tmp_path):
        from python_backend.template_engine import parse_template
        from python_backend.pptx_renderer import render_deck_with_template

        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        deck = _build_mock_deck()

        out_json = tmp_path / "output" / "deck.json"
        out_pptx = tmp_path / "output" / "deck.pptx"
        render_deck_with_template(deck, out_json, out_pptx, config)

        assert out_pptx.exists()
        assert out_json.exists()
        assert out_pptx.stat().st_size > 0

    def test_render_creates_correct_slide_count(self, tmp_path):
        from python_backend.template_engine import parse_template
        from python_backend.pptx_renderer import render_deck_with_template

        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        deck = _build_mock_deck()

        out_json = tmp_path / "output" / "deck.json"
        out_pptx = tmp_path / "output" / "deck.pptx"
        render_deck_with_template(deck, out_json, out_pptx, config)

        prs = Presentation(str(out_pptx))
        assert len(prs.slides) == 5

    def test_render_json_content(self, tmp_path):
        from python_backend.template_engine import parse_template
        from python_backend.pptx_renderer import render_deck_with_template

        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        deck = _build_mock_deck()

        out_json = tmp_path / "output" / "deck.json"
        out_pptx = tmp_path / "output" / "deck.pptx"
        render_deck_with_template(deck, out_json, out_pptx, config)

        saved_deck = json.loads(out_json.read_text())
        assert saved_deck["deckTitle"] == "Test Deck"
        assert saved_deck["slideCount"] == 5

    def test_render_all_layout_types(self, tmp_path):
        """Ensure every supported layout renders without error."""
        from python_backend.template_engine import parse_template
        from python_backend.pptx_renderer import render_deck_with_template

        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)

        all_layouts = [
            "title", "agenda", "section", "bullet", "two-column",
            "comparison", "timeline", "process", "table", "chart",
            "quote", "summary", "closing",
        ]
        slides = []
        for i, layout in enumerate(all_layouts):
            slide = {
                "page": i + 1, "layout": layout,
                "title": f"Slide {i+1}: {layout}", "subtitle": "Sub",
                "objective": "Obj", "bullets": ["Item 1", "Item 2", "Item 3"],
                "left": ["L1", "L2"], "right": ["R1", "R2"],
                "table": {"columns": ["A", "B", "C"], "rows": [["1", "2", "3"], ["4", "5", "6"]]},
                "chart": {
                    "type": "bar", "title": "Test Chart",
                    "categories": ["X", "Y", "Z"],
                    "series": [{"name": "S1", "data": [10, 20, 30]}]
                },
                "visuals": ["visual hint"], "sources": [],
                "speakerNotes": ["Note"],
            }
            slides.append(slide)

        deck = {
            "schemaVersion": "0.3.0", "deckTitle": "All Layouts",
            "language": "en", "audience": "QA", "scenario": "Test",
            "tone": "formal", "theme": "default",
            "sourceDisplayMode": "hidden", "slideCount": len(all_layouts),
            "needsSpeakerNotes": True, "assumptions": [], "slides": slides,
        }

        out_json = tmp_path / "output" / "all.json"
        out_pptx = tmp_path / "output" / "all.pptx"
        render_deck_with_template(deck, out_json, out_pptx, config)

        prs = Presentation(str(out_pptx))
        assert len(prs.slides) == len(all_layouts)

    def test_render_chart_types(self, tmp_path):
        """Test bar, line, pie, area chart rendering."""
        from python_backend.template_engine import parse_template
        from python_backend.pptx_renderer import render_deck_with_template

        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)

        for chart_type in ("bar", "line", "pie", "area"):
            slides = [{
                "page": 1, "layout": "chart", "title": f"{chart_type} chart",
                "subtitle": "", "objective": "", "bullets": ["Takeaway"],
                "left": [], "right": [],
                "table": {"columns": [], "rows": []},
                "chart": {
                    "type": chart_type, "title": f"Test {chart_type}",
                    "categories": ["A", "B", "C"],
                    "series": [{"name": "V", "data": [1, 2, 3]}]
                },
                "visuals": [], "sources": [], "speakerNotes": [],
            }]
            deck = {
                "schemaVersion": "0.3.0", "deckTitle": f"{chart_type} Test",
                "language": "en", "audience": "QA", "scenario": "Test",
                "tone": "formal", "theme": "default",
                "sourceDisplayMode": "hidden", "slideCount": 1,
                "needsSpeakerNotes": False, "assumptions": [], "slides": slides,
            }
            out_json = tmp_path / f"{chart_type}.json"
            out_pptx = tmp_path / f"{chart_type}.pptx"
            render_deck_with_template(deck, out_json, out_pptx, config)
            assert out_pptx.exists()

    def test_render_empty_chart_shows_placeholder(self, tmp_path):
        from python_backend.template_engine import parse_template
        from python_backend.pptx_renderer import render_deck_with_template

        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        slides = [{
            "page": 1, "layout": "chart", "title": "Empty Chart",
            "subtitle": "", "objective": "", "bullets": [],
            "left": [], "right": [],
            "table": {"columns": [], "rows": []},
            "chart": {"type": "bar", "title": "No Data", "categories": [], "series": []},
            "visuals": [], "sources": [], "speakerNotes": [],
        }]
        deck = {
            "schemaVersion": "0.3.0", "deckTitle": "Empty Chart",
            "language": "en", "audience": "QA", "scenario": "Test",
            "tone": "formal", "theme": "default",
            "sourceDisplayMode": "hidden", "slideCount": 1,
            "needsSpeakerNotes": False, "assumptions": [], "slides": slides,
        }
        out_json = tmp_path / "empty-chart.json"
        out_pptx = tmp_path / "empty-chart.pptx"
        render_deck_with_template(deck, out_json, out_pptx, config)
        assert out_pptx.exists()

    def test_render_table_with_data(self, tmp_path):
        from python_backend.template_engine import parse_template
        from python_backend.pptx_renderer import render_deck_with_template

        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        slides = [{
            "page": 1, "layout": "table", "title": "Data Table",
            "subtitle": "", "objective": "", "bullets": [],
            "left": [], "right": [],
            "table": {
                "columns": ["Name", "Score", "Grade"],
                "rows": [["Alice", 95, "A"], ["Bob", 82, "B"]]
            },
            "chart": {"type": "bar", "title": "", "categories": [], "series": []},
            "visuals": [], "sources": [], "speakerNotes": [],
        }]
        deck = {
            "schemaVersion": "0.3.0", "deckTitle": "Table Test",
            "language": "en", "audience": "QA", "scenario": "Test",
            "tone": "formal", "theme": "default",
            "sourceDisplayMode": "hidden", "slideCount": 1,
            "needsSpeakerNotes": False, "assumptions": [], "slides": slides,
        }
        out_json = tmp_path / "table.json"
        out_pptx = tmp_path / "table.pptx"
        render_deck_with_template(deck, out_json, out_pptx, config)
        assert out_pptx.exists()

    def test_render_with_speaker_notes(self, tmp_path):
        from python_backend.template_engine import parse_template
        from python_backend.pptx_renderer import render_deck_with_template

        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        deck = _build_mock_deck()

        out_json = tmp_path / "notes.json"
        out_pptx = tmp_path / "notes.pptx"
        render_deck_with_template(deck, out_json, out_pptx, config)

        prs = Presentation(str(out_pptx))
        # First slide should have speaker notes
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "Welcome" in notes


# ── Dual Render Path Tests ───────────────────────────────────────────


class TestDualRenderPath:
    def test_skill_api_uses_pptx_renderer_with_template(self, tmp_path):
        """When template is provided, skill_api should use python-pptx renderer."""
        from python_backend.template_engine import parse_template
        tpl = _create_test_template(tmp_path)

        deck = _build_mock_deck()
        request = {
            "action": "create",
            "prompt": "Test deck",
            "mock": True,
            "template": str(tpl),
            "_baseDir": str(tmp_path),
            "outputJson": str(tmp_path / "out.json"),
            "outputPptx": str(tmp_path / "out.pptx"),
        }

        from python_backend.skill_api import handle_skill_request
        result = handle_skill_request(request)

        assert result["ok"] is True
        assert result["renderer"] == "python-pptx"
        assert Path(result["pptxPath"]).exists()

    def test_skill_api_uses_js_renderer_without_template(self, tmp_path):
        """Without template, skill_api should use JS (pptxgenjs) renderer."""
        request = {
            "action": "create",
            "prompt": "Test deck",
            "mock": True,
            "_baseDir": str(tmp_path),
            "outputJson": str(tmp_path / "out.json"),
            "outputPptx": str(tmp_path / "out.pptx"),
        }

        from python_backend.skill_api import handle_skill_request
        result = handle_skill_request(request)

        assert result["ok"] is True
        assert result["renderer"] == "pptxgenjs"

    def test_response_includes_renderer_field(self, tmp_path):
        """Response should always include the renderer field."""
        tpl = _create_test_template(tmp_path)
        request = {
            "action": "create",
            "prompt": "Test",
            "mock": True,
            "template": str(tpl),
            "_baseDir": str(tmp_path),
            "outputJson": str(tmp_path / "out.json"),
            "outputPptx": str(tmp_path / "out.pptx"),
        }

        from python_backend.skill_api import handle_skill_request
        result = handle_skill_request(request)
        assert "renderer" in result


# ── Color Helper Tests ───────────────────────────────────────────────


class TestColorHelpers:
    def test_rgb_conversion(self):
        from python_backend.pptx_renderer import _rgb
        color = _rgb("0F766E")
        assert str(color) == "0F766E"

    def test_rgb_with_hash(self):
        from python_backend.pptx_renderer import _rgb
        color = _rgb("#2563EB")
        assert str(color) == "2563EB"

    def test_get_colors_without_template(self):
        from python_backend.pptx_renderer import _get_colors
        colors = _get_colors(None)
        assert colors["primary"] == "0F766E"

    def test_get_colors_with_template(self, tmp_path):
        from python_backend.template_engine import parse_template, ThemeColors
        from python_backend.pptx_renderer import _get_colors

        tpl = _create_test_template(tmp_path)
        config = parse_template(tpl)
        config.theme_colors = ThemeColors(primary="FF0000", secondary="00FF00", accent="0000FF")
        colors = _get_colors(config)
        assert colors["primary"] == "FF0000"
        assert colors["secondary"] == "00FF00"
