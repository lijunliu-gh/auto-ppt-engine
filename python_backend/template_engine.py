"""Template engine: parse .pptx master templates and map schema layouts to slide layouts."""

from __future__ import annotations

import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

logger = logging.getLogger("auto-ppt")

# Maximum template file size: 100 MB
MAX_TEMPLATE_SIZE = 100 * 1024 * 1024


@dataclass
class PlaceholderInfo:
    """Metadata about a single placeholder in a slide layout."""
    idx: int
    name: str
    placeholder_type: str  # e.g. "TITLE", "BODY", "SUBTITLE", "PICTURE", etc.
    left: int  # EMU
    top: int   # EMU
    width: int  # EMU
    height: int  # EMU


@dataclass
class LayoutInfo:
    """Metadata about a single slide layout in the template."""
    index: int
    name: str
    placeholders: List[PlaceholderInfo] = field(default_factory=list)

    @property
    def has_title(self) -> bool:
        return any(p.placeholder_type == "TITLE" for p in self.placeholders)

    @property
    def has_body(self) -> bool:
        return any(p.placeholder_type == "BODY" for p in self.placeholders)

    @property
    def has_picture(self) -> bool:
        return any(p.placeholder_type == "PICTURE" for p in self.placeholders)

    @property
    def body_count(self) -> int:
        return sum(1 for p in self.placeholders if p.placeholder_type == "BODY")


@dataclass
class ThemeColors:
    """Extracted theme color palette."""
    primary: str = "0F766E"     # teal (fallback)
    secondary: str = "2563EB"   # blue (fallback)
    accent: str = "F59E0B"      # amber (fallback)
    background: str = "FFFFFF"
    text_dark: str = "1E293B"
    text_light: str = "FFFFFF"


@dataclass
class TemplateConfig:
    """Parsed template configuration ready for rendering."""
    source_path: Path
    layouts: List[LayoutInfo]
    layout_mapping: Dict[str, int]  # schema layout name -> template layout index
    theme_colors: ThemeColors = field(default_factory=ThemeColors)
    font_heading: str = "Aptos Display"
    font_body: str = "Aptos"


# Default mapping from our schema layout names to common PowerPoint layout names
_LAYOUT_NAME_HINTS = {
    "title": ["title slide", "title", "封面"],
    "section": ["section header", "section", "节标题"],
    "bullet": ["title and content", "content", "标题和内容"],
    "two-column": ["two content", "comparison", "两栏内容"],
    "comparison": ["two content", "comparison", "比较"],
    "blank": ["blank", "空白"],
}


def parse_template(template_path: str | Path) -> TemplateConfig:
    """Parse a .pptx template file and extract layout/placeholder/theme information.

    Args:
        template_path: Path to a .pptx file to use as template.

    Returns:
        TemplateConfig with extracted metadata.

    Raises:
        FileNotFoundError: If template file doesn't exist.
        ValueError: If file is too large or not a valid .pptx.
        RuntimeError: If template parsing fails.
    """
    path = Path(template_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"Template not found: {path}")
    if not path.suffix.lower() == ".pptx":
        raise ValueError(f"Template must be a .pptx file, got: {path.suffix}")
    if path.stat().st_size > MAX_TEMPLATE_SIZE:
        raise ValueError(f"Template exceeds {MAX_TEMPLATE_SIZE // (1024*1024)}MB limit")

    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise RuntimeError(f"Failed to parse template: {e}") from e

    layouts = _extract_layouts(prs)
    theme_colors = _extract_theme_colors(prs)
    font_heading, font_body = _extract_fonts(prs)
    layout_mapping = _build_layout_mapping(layouts)

    config = TemplateConfig(
        source_path=path,
        layouts=layouts,
        layout_mapping=layout_mapping,
        theme_colors=theme_colors,
        font_heading=font_heading,
        font_body=font_body,
    )
    logger.info(
        "Parsed template: %s (%d layouts, %d mapped)",
        path.name, len(layouts), len(layout_mapping),
    )
    return config


def _extract_layouts(prs: Presentation) -> List[LayoutInfo]:
    """Extract all slide layouts and their placeholders."""
    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type = _placeholder_type_name(ph)
            placeholders.append(PlaceholderInfo(
                idx=ph.placeholder_format.idx,
                name=ph.name or f"Placeholder {ph.placeholder_format.idx}",
                placeholder_type=ph_type,
                left=ph.left or 0,
                top=ph.top or 0,
                width=ph.width or 0,
                height=ph.height or 0,
            ))
        layouts.append(LayoutInfo(
            index=idx,
            name=layout.name or f"Layout {idx}",
            placeholders=placeholders,
        ))
    return layouts


def _placeholder_type_name(ph: Any) -> str:
    """Get a readable placeholder type name."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        type_val = ph.placeholder_format.type
        if type_val is None:
            return "BODY"
        mapping = {
            PP_PLACEHOLDER.TITLE: "TITLE",
            PP_PLACEHOLDER.CENTER_TITLE: "TITLE",
            PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
            PP_PLACEHOLDER.BODY: "BODY",
            PP_PLACEHOLDER.OBJECT: "BODY",
        }
        return mapping.get(type_val, str(type_val).split(".")[-1].split("(")[0])
    except Exception:
        return "UNKNOWN"


def _extract_theme_colors(prs: Presentation) -> ThemeColors:
    """Extract theme colors from the presentation.

    Falls back to defaults if theme colors cannot be extracted.
    """
    colors = ThemeColors()
    try:
        theme = prs.slide_masters[0].element
        # Try to find theme XML element
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        theme_el = theme.find(".//a:theme", ns)
        if theme_el is not None:
            color_scheme = theme_el.find(".//a:clrScheme", ns)
            if color_scheme is not None:
                _parse_color_scheme(color_scheme, colors, ns)
    except Exception:
        logger.debug("Could not extract theme colors, using defaults")
    return colors


def _parse_color_scheme(color_scheme: Any, colors: ThemeColors, ns: dict) -> None:
    """Parse a clrScheme XML element into ThemeColors."""
    color_map = {
        "dk1": "text_dark",
        "lt1": "background",
        "accent1": "primary",
        "accent2": "secondary",
        "accent3": "accent",
    }
    for xml_name, attr_name in color_map.items():
        el = color_scheme.find(f"a:{xml_name}", ns)
        if el is not None:
            srgb = el.find("a:srgbClr", ns)
            if srgb is not None and "val" in srgb.attrib:
                setattr(colors, attr_name, srgb.attrib["val"])


def _extract_fonts(prs: Presentation) -> Tuple[str, str]:
    """Extract heading and body fonts from the template theme."""
    heading = "Aptos Display"
    body = "Aptos"
    try:
        theme = prs.slide_masters[0].element
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        theme_el = theme.find(".//a:theme", ns)
        if theme_el is not None:
            font_scheme = theme_el.find(".//a:fontScheme", ns)
            if font_scheme is not None:
                major = font_scheme.find("a:majorFont/a:latin", ns)
                minor = font_scheme.find("a:minorFont/a:latin", ns)
                if major is not None and "typeface" in major.attrib:
                    heading = major.attrib["typeface"]
                if minor is not None and "typeface" in minor.attrib:
                    body = minor.attrib["typeface"]
    except Exception:
        logger.debug("Could not extract fonts, using defaults")
    return heading, body


def _build_layout_mapping(layouts: List[LayoutInfo]) -> Dict[str, int]:
    """Build auto-mapping from schema layout names to template layout indices.

    Strategy:
    1. Match by layout name hints (case-insensitive substring match)
    2. For unmapped schema layouts, infer from placeholder structure
    3. Fall back to the first layout with title+body for everything else
    """
    mapping: Dict[str, int] = {}
    layout_names_lower = [(i, l.name.lower()) for i, l in enumerate(layouts)]

    # Step 1: Name-based matching
    for schema_name, hints in _LAYOUT_NAME_HINTS.items():
        for hint in hints:
            for idx, lname in layout_names_lower:
                if hint in lname and schema_name not in mapping:
                    mapping[schema_name] = idx
                    break
            if schema_name in mapping:
                break

    # Step 2: Structure-based fallbacks
    fallback_idx = _find_content_layout(layouts)

    # Map remaining schema layouts to reasonable template layouts
    all_schema_layouts = [
        "title", "agenda", "section", "bullet", "two-column",
        "comparison", "timeline", "process", "table", "chart",
        "quote", "summary", "closing",
    ]
    for name in all_schema_layouts:
        if name not in mapping:
            if name == "title":
                # Find layout with TITLE + SUBTITLE
                for l in layouts:
                    if l.has_title and any(p.placeholder_type == "SUBTITLE" for p in l.placeholders):
                        mapping[name] = l.index
                        break
            elif name == "closing":
                # Reuse title layout for closing
                if "title" in mapping:
                    mapping[name] = mapping["title"]
            elif name in ("chart", "table"):
                # Charts and tables need body area
                if fallback_idx is not None:
                    mapping[name] = fallback_idx
            if name not in mapping and fallback_idx is not None:
                mapping[name] = fallback_idx

    return mapping


def _find_content_layout(layouts: List[LayoutInfo]) -> Optional[int]:
    """Find the best general-purpose content layout (title + body)."""
    # Prefer layout with title + exactly one body placeholder
    for l in layouts:
        if l.has_title and l.body_count == 1:
            return l.index
    # Fall back to any layout with title + body
    for l in layouts:
        if l.has_title and l.has_body:
            return l.index
    # Fall back to first layout
    return 0 if layouts else None


def describe_template(config: TemplateConfig) -> Dict[str, Any]:
    """Return a human-readable summary of the parsed template for diagnostics."""
    return {
        "source": str(config.source_path),
        "layoutCount": len(config.layouts),
        "layouts": [
            {
                "index": l.index,
                "name": l.name,
                "placeholders": [
                    {"idx": p.idx, "name": p.name, "type": p.placeholder_type}
                    for p in l.placeholders
                ],
            }
            for l in config.layouts
        ],
        "layoutMapping": config.layout_mapping,
        "fonts": {"heading": config.font_heading, "body": config.font_body},
        "themeColors": {
            "primary": config.theme_colors.primary,
            "secondary": config.theme_colors.secondary,
            "accent": config.theme_colors.accent,
        },
    }


# ---------------------------------------------------------------------------
# Theme system
# ---------------------------------------------------------------------------

# CJK fallback fonts appended to any extracted font name
_CJK_FALLBACK = "Microsoft YaHei, PingFang SC, Meiryo, Noto Sans CJK SC"

# Built-in theme directory (relative to this file)
_THEMES_DIR = Path(__file__).resolve().parent.parent / "assets" / "themes"

# Default theme name when nothing is specified
DEFAULT_THEME = "business-clean"


def _font_stack(font_name: str) -> str:
    """Build a font stack with CJK fallbacks from a base font name."""
    if _CJK_FALLBACK.split(",")[0].strip() in font_name:
        return font_name  # already has CJK fallbacks
    return f"{font_name}, {_CJK_FALLBACK}"


def template_config_to_theme(config: TemplateConfig) -> Dict[str, Any]:
    """Convert a parsed TemplateConfig into a theme dict.

    This is the bridge between template_engine's extraction and
    generate-ppt.js's rendering.
    """
    tc = config.theme_colors
    return {
        "name": f"brand-{config.source_path.stem}",
        "colors": {
            "primary": tc.primary,
            "secondary": tc.secondary,
            "accent": tc.accent,
            "background": tc.background,
            "slideBg": "FFFFFF",
            "text": tc.text_dark,
            "textLight": tc.text_light,
            "textMuted": "64748B",
            "headerBg": _lighten_color(tc.primary, 0.85),
            "border": "CBD5E1",
            "closingBg": tc.text_dark,
            "titleBg": tc.background,
        },
        "fonts": {
            "heading": _font_stack(config.font_heading),
            "body": _font_stack(config.font_body),
        },
        "chartColors": [tc.primary, tc.secondary, tc.accent, "DC2626", "7C3AED", "059669"],
    }


def load_theme(theme_name: str | None = None) -> Dict[str, Any]:
    """Load a built-in theme by name.

    Args:
        theme_name: Name of a built-in theme (e.g. 'business-clean').
                    If None or not found, falls back to DEFAULT_THEME.

    Returns:
        Theme dict matching the theme-schema.json structure.
    """
    name = theme_name or DEFAULT_THEME
    theme_path = _THEMES_DIR / f"{name}.json"
    if not theme_path.exists():
        logger.warning("Theme '%s' not found, falling back to '%s'", name, DEFAULT_THEME)
        theme_path = _THEMES_DIR / f"{DEFAULT_THEME}.json"
    if not theme_path.exists():
        raise FileNotFoundError(f"Default theme not found: {theme_path}")
    with open(theme_path, "r", encoding="utf-8") as f:
        return json.load(f)


def resolve_theme(
    template_path: str | Path | None = None,
    theme_name: str | None = None,
) -> Dict[str, Any]:
    """Resolve the theme to use for rendering.

    Priority:
    1. If template_path is provided, parse it and extract theme
    2. If theme_name is provided, load the built-in theme
    3. Fall back to DEFAULT_THEME

    Returns:
        Theme dict ready to be injected into the deck JSON.
    """
    if template_path:
        config = parse_template(template_path)
        theme = template_config_to_theme(config)
        logger.info("Resolved theme from template: %s", theme["name"])
        return theme

    theme = load_theme(theme_name)
    logger.info("Resolved built-in theme: %s", theme["name"])
    return theme


def _lighten_color(hex_color: str, factor: float) -> str:
    """Lighten a hex color by blending with white.

    Args:
        hex_color: 6-digit hex color (no # prefix)
        factor: 0.0 = original, 1.0 = white
    """
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        r = int(r + (255 - r) * factor)
        g = int(g + (255 - g) * factor)
        b = int(b + (255 - b) * factor)
        return f"{r:02X}{g:02X}{b:02X}"
    except (ValueError, IndexError):
        return "E2E8F0"
